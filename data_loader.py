"""
Data loader module for processing various file formats
"""
import io
from typing import List, Dict, Any
import fitz  # PyMuPDF
import docx2txt
import pandas as pd
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from utils import get_logger, clean_text, validate_file_type

logger = get_logger(__name__)


class DataLoader:
    """Load and process various file formats"""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Initialize DataLoader
        
        Args:
            chunk_size: Size of text chunks
            chunk_overlap: Overlap between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
    
    def load_pdf(self, file_bytes: bytes, filename: str) -> str:
        """
        Load text from PDF file
        
        Args:
            file_bytes: PDF file bytes
            filename: Name of the file
            
        Returns:
            Extracted text
        """
        try:
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            logger.info(f"Successfully loaded PDF: {filename}")
            return text
        except Exception as e:
            logger.error(f"Error loading PDF {filename}: {e}")
            return ""
    
    def load_docx(self, file_bytes: bytes, filename: str) -> str:
        """
        Load text from DOCX file
        
        Args:
            file_bytes: DOCX file bytes
            filename: Name of the file
            
        Returns:
            Extracted text
        """
        try:
            text = docx2txt.process(io.BytesIO(file_bytes))
            logger.info(f"Successfully loaded DOCX: {filename}")
            return text
        except Exception as e:
            logger.error(f"Error loading DOCX {filename}: {e}")
            return ""
    
    def load_txt(self, file_bytes: bytes, filename: str) -> str:
        """
        Load text from TXT file
        
        Args:
            file_bytes: TXT file bytes
            filename: Name of the file
            
        Returns:
            Extracted text
        """
        try:
            text = file_bytes.decode('utf-8')
            logger.info(f"Successfully loaded TXT: {filename}")
            return text
        except Exception as e:
            logger.error(f"Error loading TXT {filename}: {e}")
            return ""
    
    def load_csv(self, file_bytes: bytes, filename: str) -> str:
        """
        Load text from CSV file
        
        Args:
            file_bytes: CSV file bytes
            filename: Name of the file
            
        Returns:
            Extracted text
        """
        try:
            df = pd.read_csv(io.BytesIO(file_bytes))
            # Convert DataFrame to text representation
            text = df.to_string()
            logger.info(f"Successfully loaded CSV: {filename}")
            return text
        except Exception as e:
            logger.error(f"Error loading CSV {filename}: {e}")
            return ""
    
    def load_pptx(self, file_bytes: bytes, filename: str) -> str:
        """
        Load text from PPTX file
        
        Args:
            file_bytes: PPTX file bytes
            filename: Name of the file
            
        Returns:
            Extracted text
        """
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            logger.info(f"Successfully loaded PPTX: {filename}")
            return text
        except Exception as e:
            logger.error(f"Error loading PPTX {filename}: {e}")
            return ""
    
    def load_xlsx(self, file_bytes: bytes, filename: str) -> str:
        """
        Load text from XLSX file
        
        Args:
            file_bytes: XLSX file bytes
            filename: Name of the file
            
        Returns:
            Extracted text
        """
        try:
            df = pd.read_excel(io.BytesIO(file_bytes))
            text = df.to_string()
            logger.info(f"Successfully loaded XLSX: {filename}")
            return text
        except Exception as e:
            logger.error(f"Error loading XLSX {filename}: {e}")
            return ""
    
    def load_file(self, file_bytes: bytes, filename: str) -> str:
        """
        Load file based on extension
        
        Args:
            file_bytes: File bytes
            filename: Name of the file
            
        Returns:
            Extracted text
        """
        if not validate_file_type(filename):
            logger.warning(f"Unsupported file type: {filename}")
            return ""
        
        extension = filename.lower().split('.')[-1]
        
        loaders = {
            'pdf': self.load_pdf,
            'docx': self.load_docx,
            'txt': self.load_txt,
            'csv': self.load_csv,
            'pptx': self.load_pptx,
            'xlsx': self.load_xlsx
        }
        
        loader = loaders.get(extension)
        if loader:
            return loader(file_bytes, filename)
        else:
            logger.warning(f"No loader found for extension: {extension}")
            return ""
    
    def process_files(self, files: List[tuple]) -> List[Document]:
        """
        Process multiple files and create Document objects
        
        Args:
            files: List of tuples (file_bytes, filename)
            
        Returns:
            List of Document objects with chunks
        """
        all_documents = []
        
        for file_bytes, filename in files:
            # Load text from file
            text = self.load_file(file_bytes, filename)
            
            if not text:
                logger.warning(f"No text extracted from {filename}")
                continue
            
            # Clean text
            cleaned_text = clean_text(text)
            
            # Split into chunks
            chunks = self.text_splitter.split_text(cleaned_text)
            
            # Create Document objects
            for idx, chunk in enumerate(chunks):
                doc = Document(
                    page_content=chunk,
                    metadata={
                        "source": filename,
                        "chunk_id": idx,
                        "total_chunks": len(chunks)
                    }
                )
                all_documents.append(doc)
            
            logger.info(f"Created {len(chunks)} chunks from {filename}")
        
        logger.info(f"Total documents created: {len(all_documents)}")
        return all_documents
    
    def get_file_summary(self, files: List[tuple]) -> Dict[str, Any]:
        """
        Get summary statistics of processed files
        
        Args:
            files: List of tuples (file_bytes, filename)
            
        Returns:
            Dictionary with summary statistics
        """
        summary = {
            "total_files": len(files),
            "files": [],
            "total_characters": 0,
            "total_chunks": 0
        }
        
        for file_bytes, filename in files:
            text = self.load_file(file_bytes, filename)
            if text:
                chunks = self.text_splitter.split_text(text)
                file_info = {
                    "filename": filename,
                    "characters": len(text),
                    "chunks": len(chunks)
                }
                summary["files"].append(file_info)
                summary["total_characters"] += len(text)
                summary["total_chunks"] += len(chunks)
        
        return summary